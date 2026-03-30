"""
Elementi grafici per il renderer PPTX nativo: variare studio vs presentazione senza template .pptx.
"""
from __future__ import annotations

import hashlib
import os

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from pptx.dml.color import RGBColor


def _blend(a: RGBColor, b: RGBColor, t: float) -> RGBColor:
    t = max(0.0, min(1.0, t))
    try:
        ar, ag, ab = int(a[0]), int(a[1]), int(a[2])
        br, bg, bb = int(b[0]), int(b[1]), int(b[2])
    except Exception:
        return a
    return RGBColor(
        int(ar + (br - ar) * t),
        int(ag + (bg - ag) * t),
        int(ab + (bb - ab) * t),
    )


def _tune(rgb: RGBColor, factor: float) -> RGBColor:
    try:
        r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
    except Exception:
        return rgb
    f = max(0.15, min(1.9, factor))
    return RGBColor(
        min(255, max(0, int(r * f))),
        min(255, max(0, int(g * f))),
        min(255, max(0, int(b * f))),
    )


def add_title_backdrop_shapes(
    slide,
    *,
    deck_mode: str,
    profile_key: str,
    slide_style: int,
    accent_rgb: RGBColor,
    accent2_rgb: RGBColor,
    slide_bg_rgb: RGBColor,
) -> None:
    """Forme dietro al testo del titolo (chiamare prima dei textbox)."""
    study = deck_mode == "study"
    seed = f"{profile_key}|{slide_style}|{'study' if study else 'pres'}"
    v = int(hashlib.md5(seed.encode("utf-8")).hexdigest(), 16) % 6

    if study:
        band = _blend(slide_bg_rgb, accent_rgb, 0.2)
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.36)).fill.solid()
        sb = slide.shapes[-1]
        sb.fill.fore_color.rgb = band
        sb.line.width = Pt(0)
        if v % 2 == 0:
            for i in range(4):
                slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(11.55 + (i % 2) * 0.34),
                    Inches(0.48 + (i // 2) * 0.36),
                    Inches(0.11),
                    Inches(0.11),
                ).fill.solid()
                d = slide.shapes[-1]
                d.fill.fore_color.rgb = _tune(accent_rgb, 0.72)
                d.line.width = Pt(0)
        return

    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.78), Inches(0), Inches(0.55), Inches(7.5)).fill.solid()
    edge = slide.shapes[-1]
    edge.fill.fore_color.rgb = accent_rgb
    edge.line.width = Pt(0)
    if v >= 2:
        slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.9), Inches(4.85), Inches(4.8), Inches(4.4)).fill.solid()
        ov = slide.shapes[-1]
        ov.fill.fore_color.rgb = _blend(slide_bg_rgb, accent2_rgb, 0.14)
        ov.line.width = Pt(0)


def add_content_slide_extras(
    slide,
    *,
    deck_mode: str,
    profile_key: str,
    slide_idx: int,
    slide_type: str,
    slide_style: int,
    accent_rgb: RGBColor,
    accent2_rgb: RGBColor,
    slide_bg_rgb: RGBColor,
) -> None:
    if os.getenv("ENABLE_NATIVE_PPTX_DECOR", "1").strip().lower() not in ("1", "true", "yes"):
        return
    if slide_type in ("title", "summary", "quote", "section"):
        return
    study = deck_mode == "study"
    seed = f"{profile_key}|{slide_idx}|{slide_type}|{slide_style}"
    h = int(hashlib.md5(seed.encode("utf-8")).hexdigest(), 16)

    if study:
        if slide_type in ("bullets", "numbered", "text", "split"):
            n = 3 + (h % 3)
            base_x = 11.8 + (h % 4) * 0.02
            for i in range(n):
                slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(base_x - i * 0.13),
                    Inches(0.38 + i * 0.1),
                    Inches(0.52),
                    Inches(0.07),
                ).fill.solid()
                r = slide.shapes[-1]
                r.fill.fore_color.rgb = _tune(accent_rgb, 0.42 + 0.1 * (i % 3))
                r.line.width = Pt(0)
        return

    if slide_type in ("text", "bullets", "numbered", "split") and (h % 4) != 0:
        slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.65), Inches(4.0), Inches(5.2), Inches(4.8)).fill.solid()
        ov = slide.shapes[-1]
        ov.fill.fore_color.rgb = _blend(slide_bg_rgb, accent2_rgb, 0.16)
        ov.line.width = Pt(0)

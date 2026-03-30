import os
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from playwright.sync_api import sync_playwright
import math
import hashlib
from pathlib import Path
import ast
import json
import httpx
from urllib.parse import urlparse

from .layout_profiles import (
    resolve_topic_layout,
    section_bar_is_vertical,
    slide_style_index,
    summary_is_labeled,
)

OUTPUT_DIR = "outputs"
os.makedirs(OUTPUT_DIR, exist_ok=True)

FONT = "-apple-system, BlinkMacSystemFont, 'Segoe UI', Arial, sans-serif"
TEMPLATES_DIR = Path(__file__).resolve().parents[1] / "templates" / "pptx"
TEMPLATES_CACHE_DIR = Path(__file__).resolve().parents[1] / "templates" / ".cache"
TEMPLATES_CACHE_DIR.mkdir(parents=True, exist_ok=True)

# Evita GET ripetuti al manifest R2 nella stessa richiesta / worker.
_remote_tpl_cache: list[Path] | None = None
_remote_tpl_cache_key: str | None = None


def _valid_template_path(p: Path) -> bool:
    if p.suffix.lower() != ".pptx":
        return False
    if p.name.startswith("~$") or p.name.startswith("._"):
        return False
    if "__MACOSX" in p.parts:
        return False
    return True

def _discover_templates():
    if not TEMPLATES_DIR.exists():
        return []
    return sorted([p for p in TEMPLATES_DIR.rglob("*.pptx") if _valid_template_path(p)])


def _discover_all_template_paths():
    """Template locali + file scaricati dal manifest R2 (cache), senza duplicati."""
    seen: set[str] = set()
    out: list[Path] = []
    for p in _discover_remote_templates() + _discover_templates():
        try:
            key = str(p.resolve())
        except Exception:
            key = str(p)
        if key not in seen:
            seen.add(key)
            out.append(p)
    return out

def _safe_template_filename(name: str) -> str:
    cleaned = "".join(ch for ch in (name or "template") if ch.isalnum() or ch in ("-", "_", "."))
    if not cleaned.lower().endswith(".pptx"):
        cleaned += ".pptx"
    return cleaned or "template.pptx"

def _download_remote_template(url: str, template_id: str) -> Path | None:
    if not url:
        return None
    try:
        # Nome file univoco dalla URL (stesso "id" per Presentation/Study non deve sovrascrivere la cache).
        path_part = (urlparse(url).path or "").strip("/")
        base_from_url = Path(path_part).name if path_part else ""
        if base_from_url.lower().endswith(".pptx") and base_from_url:
            local_name = _safe_template_filename(base_from_url)
        else:
            ext = ".pptx" if not url.lower().endswith(".pptx") else ""
            local_name = _safe_template_filename(f"{template_id}{ext}")
        local_path = TEMPLATES_CACHE_DIR / local_name
        if local_path.exists() and local_path.stat().st_size > 0:
            return local_path
        with httpx.Client(timeout=30.0, follow_redirects=True) as client:
            r = client.get(url)
            r.raise_for_status()
            local_path.write_bytes(r.content)
        return local_path if local_path.stat().st_size > 0 else None
    except Exception as e:
        print(f"remote template download failed ({template_id}): {e}")
        return None


def _is_from_remote_cache(p: Path) -> bool:
    """True se il .pptx è nella cartella cache del manifest (scaricato da R2)."""
    try:
        return p.resolve().parent == TEMPLATES_CACHE_DIR.resolve()
    except Exception:
        return False


def _prefer_named_template_match(paths: list[Path]) -> Path | None:
    """Se ci sono più candidati (es. locale + R2), preferisci il file scaricato dal manifest."""
    if not paths:
        return None
    remote = [p for p in paths if _is_from_remote_cache(p)]
    pool = remote if remote else paths
    return min(pool, key=lambda x: len(x.name))


def _discover_remote_templates():
    """
    Legge un manifest JSON remoto (es. Cloudflare R2 public URL) con formato:
    { "templates": [ { "id": "boxvie", "url": "https://.../boxvie.pptx" }, ... ] }
    oppure lista diretta [{...}, {...}]
    Più voci con lo stesso "id" sono ok se le URL hanno nomi file diversi (cache per nome file nell'URL).
    """
    global _remote_tpl_cache, _remote_tpl_cache_key
    manifest_url = os.getenv("TEMPLATE_MANIFEST_URL", "").strip()
    if not manifest_url:
        return []
    if _remote_tpl_cache is not None and _remote_tpl_cache_key == manifest_url:
        return list(_remote_tpl_cache)

    try:
        with httpx.Client(timeout=20.0, follow_redirects=True) as client:
            r = client.get(manifest_url)
            r.raise_for_status()
            payload = r.json()
    except Exception as e:
        print(f"template manifest fetch failed: {e}")
        return []

    items = payload.get("templates", []) if isinstance(payload, dict) else payload
    out = []
    if not isinstance(items, list):
        return out
    for i, item in enumerate(items):
        if not isinstance(item, dict):
            continue
        tid = str(item.get("id") or f"tpl_{i}")
        url = str(item.get("url") or "").strip()
        if not url:
            continue
        p = _download_remote_template(url, tid)
        if p and _valid_template_path(p):
            out.append(p)
    _remote_tpl_cache = out
    _remote_tpl_cache_key = manifest_url
    return out


def _find_named_deck_template(study: bool) -> Path | None:
    """
    Trova un .pptx il cui nome contiene il marker della modalità, es.:
    - Nome_Template_Study.pptx (sottostringa '_Template_Study')
    - oppure Template_Study.pptx (inizia per 'Template_Study')
    """
    matches = []
    for p in _discover_all_template_paths():
        n = p.name.lower()
        if study:
            ok = "_template_study" in n or n.startswith("template_study")
        else:
            ok = False
        if ok:
            matches.append(p)
    return _prefer_named_template_match(matches)


def _is_presentation_template_filename(name: str) -> bool:
    n = name.lower()
    return "_template_presentation" in n or n.startswith("template_presentation")


def _find_presentation_template(data: dict) -> Path | None:
    """
    Presentazione: distingue template lavoro vs scolastico dal nome file.
    - …_Template_Presentation_Work… — contesto lavoro/aziendale
    - …_Template_Presentation_School… — scuola/università
    - …_Template_Presentation… senza Work/School — neutro (fallback per entrambi)
    """
    audience = (data.get("presentation_audience") or "school").strip().lower()
    if audience not in ("work", "school"):
        audience = "school"

    all_p = [p for p in _discover_all_template_paths() if _is_presentation_template_filename(p.name)]
    if not all_p:
        return None

    def has_work(nm: str) -> bool:
        stem = Path(nm).stem.lower()
        return "_template_presentation" in stem and stem.endswith("_work")

    def has_school(nm: str) -> bool:
        stem = Path(nm).stem.lower()
        return "_template_presentation" in stem and stem.endswith("_school")

    neutral = [p for p in all_p if not has_work(p.name) and not has_school(p.name)]

    if audience == "work":
        w = [p for p in all_p if has_work(p.name)]
        if w:
            return _prefer_named_template_match(w)
        if neutral:
            return _prefer_named_template_match(neutral)
        return _prefer_named_template_match(all_p)

    s = [p for p in all_p if has_school(p.name)]
    if s:
        return _prefer_named_template_match(s)
    if neutral:
        return _prefer_named_template_match(neutral)
    return _prefer_named_template_match(all_p)


def _pick_template(data: dict):
    """
    Seleziona il .pptx in base a deck_mode:
    - study → '_Template_Study' / Template_Study
    - presentation → audience work/school → Template_Presentation_Work / _School, altrimenti neutro
    Poi: manifest remoto, altrimenti qualsiasi template locale (hash deterministico).
    """
    mode = (data.get("deck_mode") or "presentation").strip().lower()
    if mode == "study":
        named = _find_named_deck_template(study=True)
        if named is not None:
            return named
    else:
        named = _find_presentation_template(data)
        if named is not None:
            return named

    # Con manifest R2 configurato e download ok: il fallback non deve pescare template vecchi solo locali.
    manifest_url = os.getenv("TEMPLATE_MANIFEST_URL", "").strip()
    remote_only = _discover_remote_templates()
    if manifest_url and remote_only:
        templates = remote_only
    else:
        templates = _discover_all_template_paths()
    if not templates:
        return None
    seed = (
        (data.get("title") or "")
        + "|"
        + (data.get("subtitle") or "")
        + "|"
        + mode
        + "|"
        + str(data.get("presentation_audience") or "")
    )
    idx = int(hashlib.sha256(seed.encode("utf-8")).hexdigest(), 16) % len(templates)
    return templates[idx]

def _clear_all_slides(prs: Presentation):
    # python-pptx non espone delete slide pubblico: usiamo API interna in modo sicuro
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

def _first_body_placeholder(slide):
    for sh in slide.placeholders:
        try:
            ph_type = sh.placeholder_format.type
            # BODY(2) / OBJECT(7) / TEXT(14) / CONTENT(19)
            if int(ph_type) in (2, 7, 14, 19):
                return sh
        except Exception:
            continue
    return None


def _theme_dict(data: dict) -> dict:
    return data.get("theme") if isinstance(data.get("theme"), dict) else {}


def _theme_rgb(theme: dict, key: str, default_hex: str) -> RGBColor:
    raw = str(theme.get(key) or default_hex).lstrip("#")
    if len(raw) < 6:
        raw = default_hex.lstrip("#")
    try:
        return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
    except Exception:
        return hex_to_rgb("#" + default_hex)


def _apply_slide_background_solid(slide, rgb: RGBColor):
    """Copre lo sfondo master/template con un colore pieno (niente immagini stock di sfondo)."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _remove_shape_safely(shape):
    try:
        el = shape.element
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)
    except Exception:
        pass


# Placeholder layout "inserisci immagine" (cerchio/icona): shape_type PLACEHOLDER, non PICTURE
_PH_TYPES_TO_STRIP = frozenset(
    {
        int(PP_PLACEHOLDER.PICTURE),
        int(PP_PLACEHOLDER.BITMAP),
        int(PP_PLACEHOLDER.MEDIA_CLIP),
        int(PP_PLACEHOLDER.SLIDE_IMAGE),
        int(PP_PLACEHOLDER.ORG_CHART),
        int(PP_PLACEHOLDER.CHART),
    }
)


def _flatten_shapes(slide):
    """Tutte le forme incluso dentro i GROUP (layout spesso annidati)."""
    acc = []

    def walk(coll):
        for s in list(coll):
            acc.append(s)
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(s.shapes)

    walk(slide.shapes)
    return acc


def _remove_picture_and_media_placeholders(slide):
    """Rimuove i riquadri "immagine" / media del template (es. cerchio con icona foto)."""
    for shape in _flatten_shapes(slide):
        try:
            if not getattr(shape, "is_placeholder", False):
                continue
            ph_t = int(shape.placeholder_format.type)
            if ph_t in _PH_TYPES_TO_STRIP:
                _remove_shape_safely(shape)
        except Exception:
            continue


def _remove_large_decorative_ovals(slide, prs: Presentation, area_ratio: float = 0.06):
    """
    Ovale / rettangoli arrotondati molto grandi senza testo (spesso maschera per foto stock).
    """
    slide_area = float(prs.slide_width * prs.slide_height)
    if slide_area <= 0:
        return
    ovalish = frozenset(
        {
            MSO_AUTO_SHAPE_TYPE.OVAL,
            MSO_AUTO_SHAPE_TYPE.OVAL_CALLOUT,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            MSO_AUTO_SHAPE_TYPE.ROUND_1_RECTANGLE,
            MSO_AUTO_SHAPE_TYPE.ROUND_2_DIAG_RECTANGLE,
            MSO_AUTO_SHAPE_TYPE.ROUND_2_SAME_RECTANGLE,
        }
    )
    for shape in _flatten_shapes(slide):
        try:
            if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
            if shape.auto_shape_type not in ovalish:
                continue
            if shape.width * shape.height < slide_area * area_ratio:
                continue
            tf = getattr(shape, "text_frame", None)
            if tf is not None:
                t = (tf.text or "").strip()
                if t:
                    continue
            _remove_shape_safely(shape)
        except Exception:
            continue


def _remove_large_background_pictures(slide, prs: Presentation, cover_ratio: float = 0.08):
    """
    Rimuove immagini bitmap grandi sullo slide (stock / foto di sfondo).
    Soglia più bassa per catturare anche icone-area grandi nel template.
    """
    slide_area = float(prs.slide_width * prs.slide_height)
    if slide_area <= 0:
        return
    for shape in _flatten_shapes(slide):
        try:
            if shape.shape_type not in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
                continue
            if shape.width * shape.height < slide_area * cover_ratio:
                continue
            _remove_shape_safely(shape)
        except Exception:
            continue


def _strip_template_visual_noise(slide, prs: Presentation):
    """Ordine: placeholder immagine → ovali decorativi → bitmap grandi → poi sfondo solido (chiamato fuori)."""
    _remove_picture_and_media_placeholders(slide)
    _remove_large_decorative_ovals(slide, prs)
    _remove_large_background_pictures(slide, prs)


def _count_picture_placeholders_on_layout(layout) -> int:
    n = 0
    for shape in layout.shapes:
        try:
            if not shape.is_placeholder:
                continue
            ph_t = int(shape.placeholder_format.type)
            if ph_t in _PH_TYPES_TO_STRIP:
                n += 1
        except Exception:
            continue
    return n


def _layout_has_text_body_placeholder(layout) -> bool:
    for shape in layout.shapes:
        try:
            if not shape.is_placeholder:
                continue
            ph_t = int(shape.placeholder_format.type)
            if ph_t in (2, 7, 14, 19):
                return True
        except Exception:
            continue
    return False


def _pick_text_friendly_layout(prs: Presentation):
    """
    Non usare sempre slide_layouts[1]: spesso è una griglia con molti placeholder immagine vuoti.
    Scegli un layout con corpo testo e al massimo un placeholder foto; se solo griglie (≥2 foto), layout vuoto.
    """
    layouts = list(prs.slide_layouts)
    if not layouts:
        return None
    blank = layouts[6] if len(layouts) > 6 else layouts[0]
    candidates: list[tuple[int, object]] = []
    for lo in layouts:
        if not _layout_has_text_body_placeholder(lo):
            continue
        pics = _count_picture_placeholders_on_layout(lo)
        candidates.append((pics, lo))
    if not candidates:
        return blank
    candidates.sort(key=lambda x: (x[0], 0))
    min_pics, best = candidates[0]
    if min_pics >= 2:
        return blank
    return best


def _style_title_shape(shape, text: str, color_rgb: RGBColor, font_title: str | None, size_pt: int = 30):
    if not getattr(shape, "text_frame", None):
        return
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text or ""
    p.font.bold = True
    p.font.size = Pt(size_pt)
    p.font.color.rgb = color_rgb
    if font_title:
        p.font.name = font_title


def _normalize_slide_type(raw: str | None) -> str:
    s = (raw or "text").lower().strip()
    allowed = frozenset({"text", "bullets", "section", "quote", "split", "numbered"})
    if s in allowed:
        return s
    if s in ("bullet", "list"):
        return "bullets"
    return "text"


def _parse_stringified_list(s: str) -> list[str] | None:
    """Se il modello restituisce una stringa tipo \"['a', 'b']\" invece di JSON array, estrae gli elementi."""
    s = (s or "").strip()
    if len(s) < 4 or not s.startswith("["):
        return None
    try:
        v = ast.literal_eval(s)
        if isinstance(v, (list, tuple)):
            return [str(x).strip() for x in v if str(x).strip()]
    except (ValueError, SyntaxError, MemoryError):
        pass
    try:
        v = json.loads(s)
        if isinstance(v, list):
            return [str(x).strip() for x in v if str(x).strip()]
    except (json.JSONDecodeError, TypeError):
        pass
    return None


def _coerce_list_content(content) -> list[str]:
    """Normalizza content di bullets/numbered: lista vera o stringa serializzata."""
    if isinstance(content, list):
        return [str(x).strip() for x in content if str(x).strip()]
    if isinstance(content, str):
        parsed = _parse_stringified_list(content)
        if parsed:
            return parsed
        return [content.strip()] if content.strip() else []
    return [str(content).strip()] if content else []


def _coerce_text_slide_content(content) -> str:
    """Evita che slide 'text' mostri raw list come unica stringa."""
    if isinstance(content, list):
        return "\n\n".join(str(x).strip() for x in content if str(x).strip())
    if not isinstance(content, str):
        content = str(content or "")
    s = content.strip()
    parsed = _parse_stringified_list(s)
    if parsed:
        return "\n\n".join(parsed)
    return content


def _fill_body_numbered(
    tf,
    lines: list[str],
    *,
    start_at: int,
    text_rgb: RGBColor,
    font_body: str | None,
    body_pt: int,
):
    """Elenco numerato con offset (es. colonna destra da 5 a 8)."""
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    clean = [x.strip() for x in lines if (x or "").strip()]
    first = True
    for i, line in enumerate(clean):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"{start_at + i}. {line}"
        p.level = 0
        p.font.size = Pt(body_pt)
        p.font.color.rgb = text_rgb
        if font_body:
            p.font.name = font_body
        try:
            p.space_after = Pt(10)
            p.line_spacing = 1.2
        except Exception:
            pass


def _split_content_two_columns(content) -> tuple[str, str] | None:
    """Per type split: due stringhe affiancate."""
    if isinstance(content, str):
        parsed = _parse_stringified_list(content)
        if parsed and len(parsed) >= 2:
            return str(parsed[0]).strip(), str(parsed[1]).strip()
    if isinstance(content, list) and len(content) >= 2:
        return str(content[0]).strip(), str(content[1]).strip()
    if isinstance(content, dict):
        a = content.get("left") or content.get("a") or content.get("sinistra")
        b = content.get("right") or content.get("b") or content.get("destra")
        if a is not None and b is not None:
            return str(a).strip(), str(b).strip()
    return None


def _fill_body_paragraphs(
    tf,
    lines: list[str],
    *,
    text_rgb: RGBColor,
    font_body: str | None,
    body_pt: int = 18,
    bullet: bool = False,
    numbered: bool = False,
):
    """Riempie il text frame con paragrafi, word wrap e (opz.) elenco puntato o numerato."""
    tf.clear()
    tf.word_wrap = True
    try:
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(4)
    except Exception:
        pass
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    clean_lines = [x.strip() for x in lines if (x or "").strip()]
    first = True
    for i, line in enumerate(clean_lines):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if numbered:
            p.text = f"{i + 1}. {line}"
        elif bullet:
            p.text = f"• {line}"
        else:
            p.text = line
        p.level = 0
        p.font.size = Pt(body_pt)
        p.font.color.rgb = text_rgb
        if font_body:
            p.font.name = font_body
        try:
            p.space_after = Pt(6)
            p.line_spacing = 1.15
        except Exception:
            pass


def _fill_body_plain(tf, text: str, text_rgb: RGBColor, font_body: str | None, body_pt: int = 19):
    """Testo lungo: spezza su righe vuote per paragrafi multipli."""
    raw = str(text or "").replace("\r\n", "\n")
    if "\n\n" in raw:
        lines = [p.strip() for p in raw.split("\n\n") if p.strip()]
    else:
        lines = [ln.strip() for ln in raw.split("\n") if ln.strip()]
    if not lines:
        lines = [raw] if raw else [""]
    _fill_body_paragraphs(tf, lines, text_rgb=text_rgb, font_body=font_body, body_pt=body_pt, bullet=False)

def _generate_ppt_with_template(data: dict, user_tier: str, template_path: Path) -> str:
    prs = Presentation(str(template_path))
    # Le slide di esempio nel file template non vengono riutilizzate: si svuota il deck e si creano
    # solo slide dai layout master. Eventuali pagine superflue nel .pptx originale non compaiono in output.
    if len(prs.slides) > 0:
        _clear_all_slides(prs)

    theme = _theme_dict(data)
    bg_rgb = _theme_rgb(theme, "bg_color", "0a0a0a")
    slide_bg_rgb = _theme_rgb(theme, "slide_bg_color", "111111")
    accent_rgb = _theme_rgb(theme, "accent_color", "00D4FF")
    text_rgb = _theme_rgb(theme, "text_color", "FFFFFF")
    subtitle_rgb = _theme_rgb(theme, "subtitle_color", "A0A0B0")
    font_title = (theme.get("font_title") or "").strip() or None
    font_body = (theme.get("font_body") or "").strip() or None

    # layout fallback robusto (text_content_layout evita griglie con N placeholder immagine vuoti)
    title_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[6]
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    text_content_layout = _pick_text_friendly_layout(prs) or blank_layout

    def add_watermark(slide):
        if user_tier not in ("free", "starter"):
            return
        tx = slide.shapes.add_textbox(Inches(0.35), Inches(7.05), Inches(3.2), Inches(0.4))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "made with VoiceMint"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        if font_body:
            p.font.name = font_body

    deck_mode = (data.get("deck_mode") or "presentation").strip().lower()

    def polish_slide(slide, bg: RGBColor):
        # Studio: solo testo — rimuovi riquadri immagine stock e uniforma sfondo al tema.
        # Presentazione: conserva placeholder foto e grafica del template per l’utente.
        if deck_mode == "study":
            _strip_template_visual_noise(slide, prs)
            _apply_slide_background_solid(slide, bg)
        else:
            pass

    # title
    s0 = prs.slides.add_slide(title_layout)
    polish_slide(s0, bg_rgb)
    if getattr(s0.shapes, "title", None):
        _style_title_shape(s0.shapes.title, data.get("title") or "Presentazione", text_rgb, font_title, 32)
    else:
        t = s0.shapes.add_textbox(Inches(0.9), Inches(1.2), Inches(11.5), Inches(1.4))
        _style_title_shape(t, data.get("title") or "Presentazione", text_rgb, font_title, 32)
    body0 = _first_body_placeholder(s0)
    if body0:
        _fill_body_plain(body0.text_frame, data.get("subtitle") or "", subtitle_rgb, font_body, 17)
    else:
        b = s0.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(11.0), Inches(1.5))
        _fill_body_plain(b.text_frame, data.get("subtitle") or "", subtitle_rgb, font_body, 17)
    add_watermark(s0)

    def _default_title_and_body(s, title: str):
        if getattr(s.shapes, "title", None):
            _style_title_shape(s.shapes.title, title, accent_rgb, font_title, 26)
        else:
            tb = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7))
            _style_title_shape(tb, title, accent_rgb, font_title, 26)
        return _first_body_placeholder(s)

    for slide in data.get("slides", []):
        st = _normalize_slide_type(slide.get("type"))
        title = slide.get("title") or ""
        content = slide.get("content")

        if st in ("section", "quote", "split"):
            layout = blank_layout
        elif st in ("text", "bullets", "numbered"):
            layout = text_content_layout
        else:
            layout = blank_layout

        s = prs.slides.add_slide(layout)
        polish_slide(s, slide_bg_rgb)

        if st == "section":
            bar = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.75),
                Inches(2.05),
                Inches(2.8),
                Inches(0.09),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent_rgb
            bar.line.width = Pt(0)
            tit = s.shapes.add_textbox(Inches(0.75), Inches(2.25), Inches(11.8), Inches(1.35))
            ttf = tit.text_frame
            ttf.clear()
            tp = ttf.paragraphs[0]
            tp.text = title or "Sezione"
            tp.font.bold = True
            tp.font.size = Pt(34)
            tp.font.color.rgb = text_rgb
            if font_title:
                tp.font.name = font_title
            sub = str(content or "").strip()
            if sub:
                stb = s.shapes.add_textbox(Inches(0.75), Inches(3.75), Inches(11.5), Inches(1.1))
                _fill_body_plain(stb.text_frame, sub, subtitle_rgb, font_body, 18)

        elif st == "quote":
            if title:
                hdr = s.shapes.add_textbox(Inches(0.75), Inches(0.45), Inches(11.5), Inches(0.55))
                _style_title_shape(hdr, title, accent_rgb, font_title, 22)
            qtxt = _coerce_text_slide_content(content).strip() or "—"
            qbox = s.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11.3), Inches(5.75))
            qtf = qbox.text_frame
            qtf.clear()
            qp = qtf.paragraphs[0]
            qp.text = f"“{qtxt}”"
            qp.font.size = Pt(18)
            qp.font.italic = True
            qp.font.color.rgb = text_rgb
            qp.alignment = PP_ALIGN.CENTER
            if font_body:
                qp.font.name = font_body
            qline = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(5.4),
                Inches(6.35),
                Inches(2.5),
                Inches(0.06),
            )
            qline.fill.solid()
            qline.fill.fore_color.rgb = accent_rgb
            qline.line.width = Pt(0)

        elif st == "split":
            pair = _split_content_two_columns(content)
            body = _default_title_and_body(s, title)
            if pair:
                left, right = pair
                vbar = s.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(6.58),
                    Inches(1.35),
                    Inches(0.05),
                    Inches(6.25),
                )
                vbar.fill.solid()
                vbar.fill.fore_color.rgb = accent_rgb
                vbar.line.width = Pt(0)
                lb = s.shapes.add_textbox(Inches(0.65), Inches(1.35), Inches(5.75), Inches(6.25))
                _fill_body_plain(lb.text_frame, left, text_rgb, font_body, 16)
                rb = s.shapes.add_textbox(Inches(6.78), Inches(1.35), Inches(5.85), Inches(6.25))
                _fill_body_plain(rb.text_frame, right, text_rgb, font_body, 16)
            else:
                text = _coerce_text_slide_content(content)
                if body:
                    _fill_body_plain(body.text_frame, text, text_rgb, font_body, 19)
                else:
                    box = s.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(11.0), Inches(6.05))
                    _fill_body_plain(box.text_frame, text, text_rgb, font_body, 19)

        elif st == "numbered":
            items = _coerce_list_content(content)
            body = _default_title_and_body(s, title)
            if body:
                _fill_body_paragraphs(
                    body.text_frame,
                    items,
                    text_rgb=text_rgb,
                    font_body=font_body,
                    body_pt=17 if len(items) > 8 else 18,
                    numbered=True,
                )
            else:
                box = s.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(11.0), Inches(6.05))
                _fill_body_paragraphs(
                    box.text_frame,
                    items,
                    text_rgb=text_rgb,
                    font_body=font_body,
                    body_pt=17 if len(items) > 8 else 18,
                    numbered=True,
                )

        elif st == "bullets":
            items = _coerce_list_content(content)
            body = _default_title_and_body(s, title)
            if body:
                _fill_body_paragraphs(
                    body.text_frame,
                    items,
                    text_rgb=text_rgb,
                    font_body=font_body,
                    body_pt=17 if len(items) > 8 else 18,
                    bullet=True,
                )
            else:
                box = s.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(11.0), Inches(6.05))
                _fill_body_paragraphs(
                    box.text_frame,
                    items,
                    text_rgb=text_rgb,
                    font_body=font_body,
                    body_pt=17 if len(items) > 8 else 18,
                    bullet=True,
                )

        else:
            text = _coerce_text_slide_content(content)
            body = _default_title_and_body(s, title)
            if body:
                _fill_body_plain(body.text_frame, text, text_rgb, font_body, 19)
            else:
                box = s.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(11.0), Inches(6.05))
                _fill_body_plain(box.text_frame, text, text_rgb, font_body, 19)

        add_watermark(s)

    # summary
    ss = prs.slides.add_slide(text_content_layout)
    polish_slide(ss, slide_bg_rgb)
    summary_title = data.get("summary_title", "Riepilogo")
    summary_text = _coerce_text_slide_content(data.get("summary", ""))
    if getattr(ss.shapes, "title", None):
        _style_title_shape(ss.shapes.title, summary_title, accent_rgb, font_title, 26)
    body_s = _first_body_placeholder(ss)
    if body_s:
        _fill_body_plain(body_s.text_frame, str(summary_text or ""), text_rgb, font_body, 19)
    else:
        box = ss.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(11.0), Inches(6.05))
        _fill_body_plain(box.text_frame, str(summary_text or ""), text_rgb, font_body, 19)
    add_watermark(ss)

    filename = f"{OUTPUT_DIR}/{uuid.uuid4()}.pptx"
    prs.save(filename)
    return filename

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _rgb_tune(rgb: RGBColor, factor: float) -> RGBColor:
    """Schiarisce (factor>1) o scurisce (factor<1) un colore per pannelli / card."""
    try:
        r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
    except Exception:
        return rgb
    f = max(0.2, min(1.8, factor))
    return RGBColor(
        min(255, max(0, int(r * f))),
        min(255, max(0, int(g * f))),
        min(255, max(0, int(b * f))),
    )


def render_html_to_image(html: str, output_path: str, width: int = 1280, height: int = 720):
    with sync_playwright() as p:
        browser = p.chromium.launch()
        context = browser.new_context(
            viewport={"width": width, "height": height},
            device_scale_factor=3
        )
        page = context.new_page()
        page.set_content(html, wait_until="domcontentloaded")
        page.wait_for_timeout(300)
        page.screenshot(path=output_path, full_page=False, omit_background=False)
        browser.close()

def generate_slide_html(title: str, content, theme: dict, slide_type: str = "content") -> str:
    bg = "#" + theme.get("bg_color", "0a0a0a")
    slide_bg = "#" + theme.get("slide_bg_color", "111111")
    accent = "#" + theme.get("accent_color", "00D4FF")
    text_color = "#" + theme.get("text_color", "FFFFFF")
    subtitle_color = "#" + theme.get("subtitle_color", "A0A0B0")

    if slide_type == "title":
        return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<style>
* {{ margin:0; padding:0; box-sizing:border-box; }}
body {{
    width:1280px; height:720px; overflow:hidden;
    background: linear-gradient(135deg, {bg} 0%, {slide_bg} 100%);
    font-family: {FONT};
    display:flex; align-items:center; justify-content:center;
    position: relative;
}}
.container {{ padding: 80px; width: 100%; }}
.tag {{ font-size: 13px; letter-spacing: 4px; text-transform: uppercase; color: {accent}; margin-bottom: 24px; font-weight: 600; }}
h1 {{ font-size: 64px; font-weight: 800; color: {text_color}; line-height: 1.1; margin-bottom: 24px; max-width: 900px; }}
.subtitle {{ font-size: 22px; color: {subtitle_color}; max-width: 700px; line-height: 1.6; }}
.accent-bar {{ width: 60px; height: 5px; background: {accent}; border-radius: 3px; margin-bottom: 32px; }}
.corner-decoration {{ position: absolute; bottom: 0; right: 0; width: 300px; height: 300px; background: radial-gradient(circle, {accent}22 0%, transparent 70%); }}
</style>
</head>
<body>
<div class="container">
    <div class="tag">Presentazione</div>
    <div class="accent-bar"></div>
    <h1>{title}</h1>
    <p class="subtitle">{content}</p>
</div>
<div class="corner-decoration"></div>
</body>
</html>"""

    elif slide_type == "bullets":
        items = content if isinstance(content, list) else [content]
        bullets_html = ""
        for i, item in enumerate(items):
            bullets_html += f'<div class="bullet-item"><div class="bullet-number">{str(i+1).zfill(2)}</div><div class="bullet-text">{item}</div></div>'

        return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<style>
* {{ margin:0; padding:0; box-sizing:border-box; }}
body {{ width:1280px; height:720px; overflow:hidden; background: {bg}; font-family: {FONT}; }}
.header {{ background: linear-gradient(90deg, {slide_bg}, {bg}); padding: 32px 64px; border-bottom: 2px solid {accent}44; display: flex; align-items: center; gap: 20px; }}
.header-dot {{ width: 10px; height: 10px; background: {accent}; border-radius: 50%; }}
h2 {{ font-size: 28px; font-weight: 700; color: {accent}; letter-spacing: 1px; }}
.content {{ padding: 32px 64px; display: grid; grid-template-columns: 1fr 1fr; gap: 16px; }}
.bullet-item {{ background: {slide_bg}; border: 1px solid {accent}22; border-radius: 12px; padding: 20px 24px; display: flex; align-items: flex-start; gap: 16px; }}
.bullet-number {{ font-size: 11px; font-weight: 800; color: {accent}; letter-spacing: 2px; min-width: 28px; padding-top: 2px; }}
.bullet-text {{ font-size: 16px; color: {text_color}; line-height: 1.5; }}
</style>
</head>
<body>
<div class="header"><div class="header-dot"></div><h2>{title}</h2></div>
<div class="content">{bullets_html}</div>
</body>
</html>"""

    elif slide_type == "text":
        return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<style>
* {{ margin:0; padding:0; box-sizing:border-box; }}
body {{ width:1280px; height:720px; overflow:hidden; background: {bg}; font-family: {FONT}; display: flex; flex-direction: column; }}
.header {{ background: linear-gradient(90deg, {slide_bg}, {bg}); padding: 32px 64px; border-bottom: 2px solid {accent}44; display: flex; align-items: center; gap: 20px; }}
.header-dot {{ width: 10px; height: 10px; background: {accent}; border-radius: 50%; }}
h2 {{ font-size: 28px; font-weight: 700; color: {accent}; letter-spacing: 1px; }}
.content {{ flex: 1; padding: 48px 64px; display: flex; align-items: center; }}
.text-body {{ font-size: 20px; color: {text_color}; line-height: 1.8; max-width: 1000px; border-left: 3px solid {accent}66; padding-left: 32px; }}
</style>
</head>
<body>
<div class="header"><div class="header-dot"></div><h2>{title}</h2></div>
<div class="content"><p class="text-body">{content}</p></div>
</body>
</html>"""

    elif slide_type == "summary":
        return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<style>
* {{ margin:0; padding:0; box-sizing:border-box; }}
body {{ width:1280px; height:720px; overflow:hidden; background: linear-gradient(135deg, {bg} 0%, {slide_bg} 100%); font-family: {FONT}; display: flex; align-items: center; justify-content: center; }}
.container {{ max-width: 900px; padding: 80px; text-align: center; }}
.label {{ font-size: 12px; letter-spacing: 4px; text-transform: uppercase; color: {accent}; margin-bottom: 32px; font-weight: 600; }}
.quote-mark {{ font-size: 80px; color: {accent}44; line-height: 0.5; margin-bottom: 24px; font-family: Georgia, serif; }}
.summary-text {{ font-size: 22px; color: {text_color}; line-height: 1.8; font-weight: 400; }}
.accent-line {{ width: 40px; height: 3px; background: {accent}; border-radius: 2px; margin: 32px auto 0; }}
</style>
</head>
<body>
<div class="container">
    <div class="label">Riepilogo</div>
    <div class="quote-mark">"</div>
    <p class="summary-text">{content}</p>
    <div class="accent-line"></div>
</div>
</body>
</html>"""

    return ""


def generate_ppt(data: dict, user_tier: str = "free") -> str:
    # Template .pptx disabilitati di default: layout curati nel renderer nativo sotto.
    use_templates = os.getenv("USE_PPTX_TEMPLATES", "0").strip().lower() in ("1", "true", "yes")
    template_path = None
    if use_templates:
        template_path = _pick_template(data)
        if template_path is not None:
            src = "R2 cache" if _is_from_remote_cache(template_path) else "locale"
            print(f"pptx template: {template_path.name} ({src})")
            try:
                return _generate_ppt_with_template(data, user_tier=user_tier, template_path=template_path)
            except Exception as e:
                print(f"template fallback ({template_path.name}): {e}")

    theme = data.get("theme", {}) or {}

    # Theme colors: input viene da nlp_parser come stringhe hex senza '#'
    bg_color = "#" + theme.get("bg_color", "0a0a0a")
    slide_bg_color = "#" + theme.get("slide_bg_color", "111111")
    accent_color = "#" + theme.get("accent_color", "00D4FF")
    text_color = "#" + theme.get("text_color", "FFFFFF")
    subtitle_color = "#" + theme.get("subtitle_color", "A0A0B0")
    font_title = theme.get("font_title", "Inter")
    font_body = theme.get("font_body", "Inter")

    bg_rgb = hex_to_rgb(bg_color)
    slide_bg_rgb = hex_to_rgb(slide_bg_color)
    accent_rgb = hex_to_rgb(accent_color)
    text_rgb = hex_to_rgb(text_color)
    subtitle_rgb = hex_to_rgb(subtitle_color)
    accent2_raw = theme.get("accent_secondary")
    if accent2_raw and str(accent2_raw).strip():
        accent2_rgb = hex_to_rgb("#" + str(accent2_raw).lstrip("#"))
    else:
        accent2_rgb = _rgb_tune(accent_rgb, 0.62)

    layout = resolve_topic_layout(data)
    title_variant = int(layout.get("title_variant", 0)) % 4
    t_title_pt = Pt(int(layout.get("title_font_pt", 44)))
    t_sub_pt = Pt(int(layout.get("subtitle_font_pt", 20)))
    sec_title_pt = Pt(int(layout.get("section_title_pt", 34)))

    def set_slide_background(slide, rgb):
        # Background solido: scuro e coerente con theme
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb

    def add_watermark(slide):
        # Watermark for free trial and Starter; no watermark for Professional, Enterprise, legacy "pro"
        if user_tier not in ("free", "starter"):
            return
        tx = slide.shapes.add_textbox(Inches(0.35), Inches(7.05), Inches(3.2), Inches(0.4))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "made with VoiceMint"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        p.font.name = font_body

    def add_centered_text(slide, text, x, y, w, h, size, color, bold=False, font_name=None):
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name or font_body
        p.alignment = PP_ALIGN.CENTER
        return box

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    # layout: geometria e gerarchia da theme.layout_profile (argomento), non da hash casuale

    slide_configs = [("title", data.get("title", ""), data.get("subtitle", ""))]
    for slide in data.get("slides", []):
        slide_configs.append(
            (_normalize_slide_type(slide.get("type")), slide.get("title", ""), slide.get("content"))
        )
    slide_configs.append(("summary", data.get("summary_title", "Riepilogo"), data.get("summary", "")))

    for slide_idx, (slide_type, title, content) in enumerate(slide_configs):
        slide_style = slide_style_index(layout, slide_idx)
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background(slide, bg_rgb if slide_type == "title" else slide_bg_rgb)

        if slide_type == "title":
            if title_variant == 0:
                accent_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(2.0),
                    Inches(2.35),
                    Inches(1.6),
                    Inches(0.06),
                )
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = accent_rgb
                accent_bar.line.width = Pt(0)
                add_centered_text(
                    slide,
                    title,
                    Inches(0.8),
                    Inches(1.45),
                    Inches(11.7),
                    Inches(1.5),
                    t_title_pt,
                    text_rgb,
                    bold=True,
                    font_name=font_title,
                )
                add_centered_text(
                    slide,
                    content or "",
                    Inches(2.0),
                    Inches(2.85),
                    Inches(9.3),
                    Inches(0.9),
                    t_sub_pt,
                    subtitle_rgb,
                    bold=False,
                    font_name=font_body,
                )
            elif title_variant == 1:
                rail = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.42),
                    Inches(0.45),
                    Inches(0.16),
                    Inches(6.6),
                )
                rail.fill.solid()
                rail.fill.fore_color.rgb = accent_rgb
                rail.line.width = Pt(0)
                tb = slide.shapes.add_textbox(Inches(0.85), Inches(1.25), Inches(11.5), Inches(1.65))
                tf = tb.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = title or ""
                p.font.size = t_title_pt
                p.font.bold = True
                p.font.color.rgb = text_rgb
                p.font.name = font_title
                p.alignment = PP_ALIGN.LEFT
                if content:
                    tb2 = slide.shapes.add_textbox(Inches(0.85), Inches(3.05), Inches(10.8), Inches(1.15))
                    tf2 = tb2.text_frame
                    tf2.clear()
                    p2 = tf2.paragraphs[0]
                    p2.text = str(content)
                    p2.font.size = t_sub_pt
                    p2.font.color.rgb = subtitle_rgb
                    p2.font.name = font_body
                    p2.alignment = PP_ALIGN.LEFT
            elif title_variant == 2:
                top_strip = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0),
                    Inches(0.38),
                    Inches(13.33),
                    Inches(0.09),
                )
                top_strip.fill.solid()
                top_strip.fill.fore_color.rgb = accent_rgb
                top_strip.line.width = Pt(0)
                add_centered_text(
                    slide,
                    title,
                    Inches(0.7),
                    Inches(1.55),
                    Inches(12.0),
                    Inches(1.55),
                    t_title_pt,
                    text_rgb,
                    bold=True,
                    font_name=font_title,
                )
                add_centered_text(
                    slide,
                    content or "",
                    Inches(1.2),
                    Inches(3.25),
                    Inches(10.9),
                    Inches(1.0),
                    t_sub_pt,
                    subtitle_rgb,
                    bold=False,
                    font_name=font_body,
                )
            else:
                bottom_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(3.8),
                    Inches(5.85),
                    Inches(5.7),
                    Inches(0.12),
                )
                bottom_bar.fill.solid()
                bottom_bar.fill.fore_color.rgb = accent_rgb
                bottom_bar.line.width = Pt(0)
                add_centered_text(
                    slide,
                    title,
                    Inches(0.8),
                    Inches(1.35),
                    Inches(11.7),
                    Inches(1.55),
                    t_title_pt,
                    text_rgb,
                    bold=True,
                    font_name=font_title,
                )
                add_centered_text(
                    slide,
                    content or "",
                    Inches(1.5),
                    Inches(2.95),
                    Inches(10.3),
                    Inches(1.0),
                    t_sub_pt,
                    subtitle_rgb,
                    bold=False,
                    font_name=font_body,
                )

        elif slide_type == "section":
            vert = section_bar_is_vertical(layout, slide_style)
            if vert:
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.75),
                    Inches(1.85),
                    Inches(0.12),
                    Inches(3.4),
                )
            else:
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.75),
                    Inches(2.05),
                    Inches(2.8),
                    Inches(0.09),
                )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent_rgb
            bar.line.width = Pt(0)
            ty = 2.05 if vert else 2.25
            tbox = slide.shapes.add_textbox(Inches(1.05), Inches(ty), Inches(11.5), Inches(1.45))
            tf = tbox.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title or "Sezione"
            p.font.bold = True
            p.font.size = Pt(int(layout.get("section_title_pt", 34)) + (2 if slide_style in (1, 3) else 0))
            p.font.color.rgb = text_rgb
            p.font.name = font_title
            sub = str(content or "").strip()
            if sub:
                sy = 3.65 if vert else 3.85
                sbox = slide.shapes.add_textbox(Inches(1.05), Inches(sy), Inches(11.2), Inches(1.15))
                _fill_body_plain(sbox.text_frame, sub, subtitle_rgb, font_body, 18)

        elif slide_type == "quote":
            if title:
                hdr = slide.shapes.add_textbox(Inches(0.75), Inches(0.45), Inches(11.5), Inches(0.55))
                tf = hdr.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = title
                p.font.bold = True
                p.font.size = Pt(22)
                p.font.color.rgb = accent_rgb
                p.font.name = font_title
            qtxt = _coerce_text_slide_content(content).strip() or "—"
            qbox = slide.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(11.3), Inches(4.6))
            qtf = qbox.text_frame
            qtf.clear()
            qp = qtf.paragraphs[0]
            qp.text = f"“{qtxt}”"
            qp.font.size = Pt(21)
            qp.font.italic = True
            qp.font.color.rgb = text_rgb
            qp.alignment = PP_ALIGN.CENTER
            qp.font.name = font_body
            qline = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(5.4),
                Inches(6.35),
                Inches(2.5),
                Inches(0.06),
            )
            qline.fill.solid()
            qline.fill.fore_color.rgb = accent_rgb
            qline.line.width = Pt(0)

        elif slide_type == "split":
            pair = _split_content_two_columns(content)
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
            tf = title_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title or ""
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = accent_rgb
            p.font.name = font_title
            p.alignment = PP_ALIGN.LEFT
            if pair:
                left, right = pair
                vbar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(6.58),
                    Inches(1.35),
                    Inches(0.05),
                    Inches(5.85),
                )
                vbar.fill.solid()
                vbar.fill.fore_color.rgb = accent_rgb
                vbar.line.width = Pt(0)
                lb = slide.shapes.add_textbox(Inches(0.65), Inches(1.35), Inches(5.75), Inches(5.85))
                _fill_body_plain(lb.text_frame, left, text_rgb, font_body, 17)
                rb = slide.shapes.add_textbox(Inches(6.78), Inches(1.35), Inches(5.85), Inches(5.85))
                _fill_body_plain(rb.text_frame, right, text_rgb, font_body, 17)
            else:
                body_box = slide.shapes.add_textbox(Inches(1.15), Inches(1.5), Inches(11.9), Inches(5.8))
                _fill_body_plain(
                    body_box.text_frame,
                    _coerce_text_slide_content(content),
                    text_rgb,
                    font_body,
                    19,
                )

        elif slide_type == "numbered":
            bullets = _coerce_list_content(content)[:12]
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
            tf = title_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title or ""
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = accent_rgb
            p.font.name = font_title
            p.alignment = PP_ALIGN.LEFT
            n = len(bullets)
            if not n:
                pass
            else:
                columns = 2 if n >= 8 else 1
                body_pt = 17 if n > 6 else 18
                if columns == 1:
                    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.32), Inches(11.5), Inches(5.85))
                    _fill_body_numbered(
                        box.text_frame,
                        bullets,
                        start_at=1,
                        text_rgb=text_rgb,
                        font_body=font_body,
                        body_pt=body_pt,
                    )
                else:
                    mid = (n + 1) // 2
                    left, right = bullets[:mid], bullets[mid:]
                    lb = slide.shapes.add_textbox(Inches(0.62), Inches(1.32), Inches(5.85), Inches(5.85))
                    _fill_body_numbered(
                        lb.text_frame,
                        left,
                        start_at=1,
                        text_rgb=text_rgb,
                        font_body=font_body,
                        body_pt=body_pt,
                    )
                    rb = slide.shapes.add_textbox(Inches(6.72), Inches(1.32), Inches(5.95), Inches(5.85))
                    _fill_body_numbered(
                        rb.text_frame,
                        right,
                        start_at=len(left) + 1,
                        text_rgb=text_rgb,
                        font_body=font_body,
                        body_pt=body_pt,
                    )

        elif slide_type == "bullets":
            bullets = _coerce_list_content(content)[:12]

            title_x = (
                1.1
                if layout.get("profile_key")
                in ("executive_premium", "startup_pitch", "history_editorial")
                else 0.8
            )
            title_box = slide.shapes.add_textbox(Inches(title_x), Inches(0.5), Inches(11.7), Inches(0.6))
            tf = title_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title or ""
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = accent_rgb
            p.font.name = font_title
            p.alignment = PP_ALIGN.LEFT

            n = len(bullets)
            if not n:
                pass
            else:
                columns = 2 if n >= 8 else 1
                body_pt = 17 if n > 6 else 18
                if columns == 1:
                    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.32), Inches(11.5), Inches(5.85))
                    _fill_body_paragraphs(
                        box.text_frame,
                        bullets,
                        text_rgb=text_rgb,
                        font_body=font_body,
                        body_pt=body_pt,
                        bullet=True,
                    )
                else:
                    mid = (n + 1) // 2
                    left, right = bullets[:mid], bullets[mid:]
                    lb = slide.shapes.add_textbox(Inches(0.62), Inches(1.32), Inches(5.85), Inches(5.85))
                    _fill_body_paragraphs(
                        lb.text_frame,
                        left,
                        text_rgb=text_rgb,
                        font_body=font_body,
                        body_pt=body_pt,
                        bullet=True,
                    )
                    rb = slide.shapes.add_textbox(Inches(6.72), Inches(1.32), Inches(5.95), Inches(5.85))
                    _fill_body_paragraphs(
                        rb.text_frame,
                        right,
                        text_rgb=text_rgb,
                        font_body=font_body,
                        body_pt=body_pt,
                        bullet=True,
                    )

        elif slide_type == "text":
            # Title at top
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
            tf = title_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title or ""
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = accent_rgb
            p.font.name = font_title
            p.alignment = PP_ALIGN.LEFT

            # Solo bordo verticale: la fascia orizzontale piena creava artefatti (linea nel testo)
            border = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.95),
                Inches(1.35),
                Inches(0.06),
                Inches(5.8),
            )
            border.fill.solid()
            border.fill.fore_color.rgb = accent_rgb
            border.line.color.rgb = accent_rgb
            border.line.width = Pt(0.0)
            body_left = Inches(1.15)
            body_top = Inches(1.35)

            body_box = slide.shapes.add_textbox(body_left, body_top, Inches(11.9), Inches(5.75))
            tf = body_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = None

            # Keep it readable (wrap by slide width)
            p = tf.paragraphs[0]
            p.text = _coerce_text_slide_content(content)
            p.font.size = Pt(20)
            p.font.color.rgb = text_rgb
            p.font.name = font_body
            p.alignment = PP_ALIGN.LEFT

        elif slide_type == "summary":
            summary_text = _coerce_text_slide_content(content)
            if summary_is_labeled(layout, slide_style):
                lbl = slide.shapes.add_textbox(Inches(1.0), Inches(1.15), Inches(11.3), Inches(0.45))
                ltf = lbl.text_frame
                ltf.clear()
                lp = ltf.paragraphs[0]
                lp.text = (title or "Riepilogo").upper()
                lp.font.size = Pt(11)
                lp.font.bold = True
                lp.font.color.rgb = accent_rgb
                lp.font.name = font_title
                lp.alignment = PP_ALIGN.CENTER
                quote_top = Inches(2.05)
            else:
                quote_top = Inches(2.35)

            add_centered_text(
                slide,
                f"“{summary_text}”",
                Inches(1.0),
                quote_top,
                Inches(11.3),
                Inches(3.0),
                Pt(24 if slide_style in (1, 3) else 23),
                text_rgb,
                bold=False,
                font_name=font_body,
            )

            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(5.15 if slide_style in (1, 3) else 5.35),
                Inches(5.35 if slide_style in (1, 3) else 5.25),
                Inches(3.0 if slide_style in (1, 3) else 2.6),
                Inches(0.05),
            )
            line.fill.solid()
            line.fill.fore_color.rgb = accent_rgb
            line.line.color.rgb = accent_rgb
            line.line.width = Pt(0.0)

        else:
            # Fallback: treat as text slide
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
            tf = title_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title or ""
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = accent_rgb
            p.font.name = font_title
            p.alignment = PP_ALIGN.LEFT

            body_box = slide.shapes.add_textbox(Inches(1.15), Inches(1.5), Inches(11.9), Inches(5.8))
            tf = body_box.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = _coerce_text_slide_content(content)
            p.font.size = Pt(20)
            p.font.color.rgb = text_rgb
            p.font.name = font_body
            p.alignment = PP_ALIGN.LEFT

        add_watermark(slide)

    filename = f"{OUTPUT_DIR}/{uuid.uuid4()}.pptx"
    prs.save(filename)
    return filename